from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os


def create_light_template():
    prs = Presentation()

    # Create title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

    title = slide.shapes.title
    title.text = "Title"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
    title.text_frame.paragraphs[0].font.size = Pt(44)

    subtitle = slide.placeholders[1]
    subtitle.text = "Subtitle"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)

    # Create content slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

    title = slide.shapes.title
    title.text = "Slide Title"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
    title.text_frame.paragraphs[0].font.size = Pt(36)

    content = slide.placeholders[1]
    content.text = "Content goes here"
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
    content.text_frame.paragraphs[0].font.size = Pt(18)

    # Save
    os.makedirs('frame', exist_ok=True)
    prs.save(os.path.join('frame', 'light.pptx'))


# Create light template only
create_light_template()

print("Light template created successfully!")