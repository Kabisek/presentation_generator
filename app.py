from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os
import uuid

app = Flask(__name__)

# Define template paths
TEMPLATE_PATHS = {
    'light': 'frame/light.pptx',
    'normal': 'frame/normal.pptx',
    'dark': 'frame/dark.pptx'
}


def generate_content_for_slides(topic, num_slides):
    """
    Simulate an ML model (e.g., LLaMA) to generate varied content for slides based on the topic.
    Returns a list of content points, one for each content slide.
    """
    # Base content provided by the user
    base_content = topic.strip()

    # If the user provided specific content, use it as a starting point
    if not base_content:
        base_content = "Introduction to Python basics."

    # Simulated ML-generated content based on the topic "introduction to python"
    if "python" in base_content.lower():
        ml_generated_content = [
            "Python is a high-level, interpreted programming language known for its readability.",
            "Example: Print 'Hello, World!' using print('Hello, World!').",
            "Python supports multiple data types like lists, dictionaries, and tuples.",
            "Example: Create a list with my_list = [1, 2, 3] and append with my_list.append(4).",
            "Python's simplicity makes it ideal for beginners and professionals alike."
        ]
    else:
        # Fallback for generic topics
        ml_generated_content = [f"Overview of {base_content}."] * num_slides

    # Ensure we have enough content for all slides
    if len(ml_generated_content) < num_slides:
        ml_generated_content = (ml_generated_content * (num_slides // len(ml_generated_content) + 1))[:num_slides]
    elif len(ml_generated_content) > num_slides:
        ml_generated_content = ml_generated_content[:num_slides]

    return ml_generated_content


def find_placeholder(slide, placeholder_type="content"):
    """
    Dynamically find the appropriate placeholder in a slide.
    - For 'subtitle', find the first non-title placeholder.
    - For 'content', find the first non-title placeholder.
    Returns the placeholder shape or None if not found.
    """
    title_shape = slide.shapes.title
    for shape in slide.placeholders:
        if shape == title_shape:
            continue
        if placeholder_type == "subtitle" and shape.placeholder_format.type in [2, 7]:  # Title (2), Subtitle (7)
            return shape
        if placeholder_type == "content" and shape.placeholder_format.type in [1, 4]:  # Body (1), Text (4)
            return shape
    return None


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/generate', methods=['POST'])
def generate_ppt():
    # Get form data
    title = request.form['title']
    author = request.form['author']
    content_input = request.form['content'].split('\n')  # Split content by newlines
    content_input = [line.strip() for line in content_input if line.strip()]  # Remove empty lines
    num_slides = int(request.form['num_slides'])
    template = request.form['template']

    # Validate template
    if template not in TEMPLATE_PATHS:
        return "Invalid template selected", 400

    # Load the selected template
    prs = Presentation(TEMPLATE_PATHS[template])

    # Clear any existing slides in the template to avoid unexpected content
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # Create title slide using the template's title slide layout
    title_slide_layout = prs.slide_layouts[0]  # Assuming layout 0 is title slide
    slide = prs.slides.add_slide(title_slide_layout)

    # Set title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Find the subtitle placeholder dynamically
    subtitle_placeholder = find_placeholder(slide, placeholder_type="subtitle")
    if subtitle_placeholder:
        subtitle_placeholder.text = f"By {author}"
        subtitle_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    else:
        # Fallback: Add a text box if no subtitle placeholder is found
        left = Pt(50)
        top = Pt(300)
        width = Pt(600)
        height = Pt(50)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = f"By {author}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)

    # Calculate number of content slides (total slides - 1 for title)
    num_content_slides = max(0, num_slides - 1)

    if num_content_slides > 0:
        # Generate varied content using simulated ML model
        content = generate_content_for_slides(content_input[0] if content_input else title, num_content_slides)

        # Distribute content across slides (one point per slide)
        content_slide_layout = prs.slide_layouts[1]  # Assuming layout 1 is content slide
        for i in range(num_content_slides):
            slide = prs.slides.add_slide(content_slide_layout)

            # Set slide title
            title_shape = slide.shapes.title
            title_shape.text = f"Slide {i + 1}"
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

            # Find the content placeholder dynamically
            body_shape = find_placeholder(slide, placeholder_type="content")
            if body_shape:
                text_frame = body_shape.text_frame
                text_frame.clear()  # Clear any existing text
                content_line = content[i]
                p = text_frame.add_paragraph()
                p.text = content_line
                p.font.size = Pt(18)
                p.alignment = PP_ALIGN.LEFT
                p.level = 0  # Ensure no bullet points
            else:
                # Fallback: Add a text box if no content placeholder is found
                left = Pt(50)
                top = Pt(150)
                width = Pt(600)
                height = Pt(300)
                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                p = text_frame.add_paragraph()
                p.text = content[i]
                p.alignment = PP_ALIGN.LEFT
                p.font.size = Pt(18)

    # Save the presentation
    output_filename = f"generated_ppt_{uuid.uuid4().hex}.pptx"
    output_path = os.path.join('generated', output_filename)

    # Ensure generated directory exists
    os.makedirs('generated', exist_ok=True)

    prs.save(output_path)

    # Send the file to the user
    return send_file(output_path, as_attachment=True, download_name=f"{title}.pptx")


if __name__ == '__main__':
    app.run(debug=True)